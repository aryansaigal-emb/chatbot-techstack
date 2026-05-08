# ================================================================
#  EMB RAG Chatbot — Backend
#  FastAPI + FAISS + sentence-transformers + OpenRouter + Supabase Auth
# ================================================================

import os
import io
import json
import re
import unicodedata
import asyncio
from typing import Any, List
from datetime import datetime, timedelta, timezone
from dotenv import load_dotenv
import numpy as np
import faiss
from sentence_transformers import SentenceTransformer
import requests
from supabase import create_client, Client
from fastapi import FastAPI, UploadFile, File, HTTPException, Depends
from fastapi.middleware.cors import CORSMiddleware
from fastapi.security import HTTPBearer, HTTPAuthorizationCredentials
from pydantic import BaseModel, Field
import secrets

# ── Load environment variables ────────────────────────────────────
load_dotenv(os.path.join(os.path.dirname(__file__), ".env"), override=True)
OPENROUTER_API_KEY = os.getenv("OPENROUTER_API_KEY")
OPENROUTER_MODEL = os.getenv("OPENROUTER_MODEL", "openrouter/free")
OPENROUTER_SITE_URL = os.getenv("OPENROUTER_SITE_URL")
OPENROUTER_APP_NAME = os.getenv("OPENROUTER_APP_NAME", "EMB RAG Chatbot")
OPENROUTER_DATETIME_TIMEZONE = os.getenv("OPENROUTER_DATETIME_TIMEZONE", "Asia/Kolkata")
MCP_SERVERS_JSON = os.getenv("MCP_SERVERS_JSON", "")
WORKSPACE_DIR = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
MCP_SERVERS_FILE = os.getenv(
    "MCP_SERVERS_FILE",
    os.path.join(os.path.dirname(__file__), "mcp_servers.json"),
)
MCP_ALLOWED_TOOLS = {
    tool.strip()
    for tool in os.getenv(
        "MCP_ALLOWED_TOOLS",
        (
            "read_file,read_text_file,read_media_file,read_multiple_files,"
            "list_directory,list_directory_with_sizes,directory_tree,"
            "search_files,get_file_info,list_allowed_directories"
        ),
    ).split(",")
    if tool.strip()
}
TESSERACT_CMD = os.getenv("TESSERACT_CMD", r"C:\Program Files\Tesseract-OCR\tesseract.exe")
SUPABASE_URL = os.getenv("SUPABASE_URL")
SUPABASE_KEY = os.getenv("SUPABASE_KEY")
FRONTEND_ORIGINS = [
    origin.strip()
    for origin in os.getenv("FRONTEND_ORIGINS", "http://localhost:5173").split(",")
    if origin.strip()
]

if not OPENROUTER_API_KEY:
    raise RuntimeError("OPENROUTER_API_KEY not found in .env")
if not SUPABASE_URL:
    raise RuntimeError("SUPABASE_URL not found in .env")
if not SUPABASE_KEY:
    raise RuntimeError("SUPABASE_KEY not found in .env")

# ── OpenRouter client ─────────────────────────────────────────────
OPENROUTER_CHAT_URL = "https://openrouter.ai/api/v1/chat/completions"

openrouter_headers = {
    "Authorization": f"Bearer {OPENROUTER_API_KEY}",
    "Content-Type": "application/json",
    "X-Title": OPENROUTER_APP_NAME,
}
if OPENROUTER_SITE_URL:
    openrouter_headers["HTTP-Referer"] = OPENROUTER_SITE_URL

WEB_SEARCH_PLUGIN = {
    "id": "web",
    "engine": "exa",
    "max_results": 2,
    "search_prompt": (
        "A web search was conducted for the user request. "
        "Use the results to answer with current information. "
        "Do not expose raw search snippets or tool calls."
    ),
}

DATE_TIME_KEYWORDS = (
    "date",
    "time",
    "today",
    "tomorrow",
    "yesterday",
    "day",
    "month",
    "year",
    "timezone",
    "current time",
    "current date",
)

WEB_SEARCH_KEYWORDS = (
    "latest",
    "recent",
    "news",
    "search",
    "web",
    "internet",
    "online",
    "live",
    "current",
    "today's",
    "today ",
    "now",
)

DIAGRAM_KEYWORDS = (
    "mermaid",
    "diagram",
    "flowchart",
    "sequence diagram",
    "class diagram",
    "state diagram",
    "er diagram",
    "entity relationship",
    "gantt",
    "pie chart",
    "mindmap",
    "timeline",
    "user journey",
    "gitgraph",
    "quadrant chart",
    "requirement diagram",
    "c4 diagram",
    "block diagram",
    "architecture diagram",
    "sankey",
    "xy chart",
    "packet diagram",
)


def wants_datetime_tool(text: str) -> bool:
    if wants_mermaid_diagram(text):
        return False

    lowered = text.lower().strip()
    explicit_patterns = (
        r"\bwhat(?:'s| is)?\s+(?:the\s+)?(?:current\s+)?(?:date|time)\b",
        r"\b(?:current|today'?s)\s+(?:date|time)\b",
        r"\bwhat\s+day\s+is\s+(?:it|today)\b",
        r"\bwhat\s+(?:is\s+)?today'?s\s+date\b",
        r"\btell\s+me\s+(?:the\s+)?(?:date|time)\b",
    )
    if any(re.search(pattern, lowered) for pattern in explicit_patterns):
        return True

    return lowered in {
        "date",
        "time",
        "today",
        "tomorrow",
        "yesterday",
        "current date",
        "current time",
    }


def wants_web_search_tool(text: str) -> bool:
    lowered = f" {text.lower()} "
    if wants_datetime_tool(text) and not any(word in lowered for word in (" latest", " news", " search", " web", " internet", " recent", " live")):
        return False
    return any(keyword in lowered for keyword in WEB_SEARCH_KEYWORDS)


def wants_mermaid_diagram(text: str) -> bool:
    lowered = text.lower()
    return any(keyword in lowered for keyword in DIAGRAM_KEYWORDS)


def mermaid_request_count(text: str) -> int:
    lowered = text.lower()
    if re.search(r"\b(two|2)\b", lowered):
        return 2
    if re.search(r"\b(three|3)\b", lowered):
        return 3
    return 1


def requested_mermaid_types(text: str) -> List[str]:
    lowered = text.lower()
    type_checks = [
        ("flowchart", ("flowchart", "flow chart", "process diagram")),
        ("sequence", ("sequence", "sequence diagram")),
        ("class", ("class diagram",)),
        ("state", ("state diagram", "state machine")),
        ("er", ("er diagram", "entity relationship", "database diagram")),
        ("gantt", ("gantt",)),
        ("pie", ("pie chart", "pie diagram")),
        ("mindmap", ("mindmap", "mind map")),
        ("timeline", ("timeline",)),
        ("journey", ("journey", "user journey")),
        ("gitgraph", ("gitgraph", "git graph")),
        ("quadrant", ("quadrant", "priority matrix")),
        ("requirement", ("requirement", "requirements diagram")),
        ("sankey", ("sankey",)),
        ("xy", ("xy chart", "line chart")),
    ]
    found = []
    for diagram_type, needles in type_checks:
        if any(needle in lowered for needle in needles):
            found.append(diagram_type)
    return found or ["flowchart"]


def mermaid_title_for(text: str, diagram_type: str) -> str:
    lowered = text.lower()
    if diagram_type == "pie" and (
        "average person" in lowered or "spends their day" in lowered or "average day" in lowered
    ):
        return "Average Day"
    if diagram_type == "pie" and ("project" in lowered or "work distribution" in lowered):
        return "Project Work Distribution"
    if "umbrella" in lowered or "rain" in lowered:
        return "Umbrella Decision"
    if "morning" in lowered:
        return "Morning Routine"
    if "average person" in lowered or "spends their day" in lowered:
        return "Average Day"
    if "login" in lowered or "upload" in lowered or "rag" in lowered:
        return "Chatbot Workflow"
    if "users" in lowered and ("messages" in lowered or "sessions" in lowered):
        return "Chat Memory Data Model"
    return diagram_type.replace("-", " ").title()


def mermaid_code_for(text: str, diagram_type: str) -> str:
    lowered = text.lower()
    title = mermaid_title_for(text, diagram_type)

    if diagram_type == "flowchart":
        if "umbrella" in lowered or "rain" in lowered:
            return """flowchart TD
    A[Check weather] --> B{Is it raining now?}
    B -->|Yes| C[Carry an umbrella]
    B -->|No| D{Is there a chance of rain?}
    D -->|High chance| C
    D -->|Low chance| E[No umbrella needed]
    C --> F[Leave prepared]
    E --> F"""
        if "morning" in lowered:
            return """flowchart TD
    A[Wake up] --> B[Brush teeth]
    B --> C[Exercise]
    C --> D[Shower]
    D --> E[Breakfast]
    E --> F[Plan the day]
    F --> G[Start work or school]"""
        if "login" in lowered or "upload" in lowered or "rag" in lowered:
            return """flowchart TD
    A[User logs in] --> B[Upload file]
    B --> C[Extract text]
    C --> D[Create chunks]
    D --> E[Embed chunks]
    E --> F[Store in vector index]
    F --> G[Ask question]
    G --> H[Retrieve relevant chunks]
    H --> I[Generate answer]
    I --> J[Show response]"""
        return """flowchart TD
    A[Start] --> B[Collect input]
    B --> C{Decision needed?}
    C -->|Yes| D[Evaluate options]
    C -->|No| E[Continue process]
    D --> F[Choose next step]
    E --> F
    F --> G[Finish]"""

    if diagram_type == "pie":
        if "average person" in lowered or "spends their day" in lowered or "day" in lowered:
            return """pie title Average Day
    "Sleep" : 8
    "Work or school" : 8
    "Leisure" : 3
    "Meals" : 2
    "Commute" : 1
    "Exercise" : 1
    "Personal tasks" : 1"""
        if "project" in lowered or "work distribution" in lowered:
            return """pie title Project Work Distribution
    "Backend" : 35
    "Frontend" : 30
    "Testing" : 20
    "Deployment" : 15"""
        return f"""pie title {title}
    "Category A" : 40
    "Category B" : 30
    "Category C" : 20
    "Other" : 10"""

    if diagram_type == "sequence":
        return """sequenceDiagram
    autonumber
    actor User
    participant Frontend
    participant Backend
    participant VectorStore
    participant Model
    User->>Frontend: Submit request
    Frontend->>Backend: Send API call
    Backend->>VectorStore: Retrieve context
    VectorStore-->>Backend: Return chunks
    Backend->>Model: Generate answer
    Model-->>Backend: Return response
    Backend-->>Frontend: Send answer
    Frontend-->>User: Render result"""

    if diagram_type == "class":
        return """classDiagram
    class User {
      +string user_id
      +login()
      +logout()
    }
    class ChatSession {
      +string id
      +string title
      +Message[] messages
    }
    class Message {
      +string role
      +string content
    }
    User "1" --> "*" ChatSession
    ChatSession "1" --> "*" Message"""

    if diagram_type == "state":
        return """stateDiagram-v2
    [*] --> Idle
    Idle --> Uploading: user uploads file
    Uploading --> Ready: indexing complete
    Ready --> Answering: user asks question
    Answering --> Ready: answer shown
    Uploading --> Error: upload fails
    Answering --> Error: model fails
    Error --> Ready: retry"""

    if diagram_type == "er":
        return """erDiagram
    USER ||--o{ CHAT_SESSION : owns
    CHAT_SESSION ||--o{ MESSAGE : contains
    USER {
      string user_id
      string passcode
      datetime created_at
    }
    CHAT_SESSION {
      string id
      string title
      datetime updated_at
    }
    MESSAGE {
      string role
      string content
    }"""

    if diagram_type == "gantt":
        return """gantt
    title Project Plan
    dateFormat  YYYY-MM-DD
    section Backend
    API routes        :a1, 2026-05-08, 1d
    Integrations      :a2, after a1, 2d
    section Frontend
    UI rendering      :b1, 2026-05-08, 2d
    Testing           :b2, after b1, 1d"""

    if diagram_type == "mindmap":
        return """mindmap
  root((Chatbot))
    Inputs
      Documents
      Questions
      Diagram prompts
    Processing
      RAG
      MCP
      Mermaid
    Outputs
      Answers
      Tables
      Diagrams"""

    if diagram_type == "timeline":
        return """timeline
    title Chatbot Request Flow
    Login : User signs in
    Upload : File is processed
    Ask : User sends prompt
    Retrieve : Context is found
    Render : Answer or diagram is shown"""

    if diagram_type == "journey":
        return """journey
    title User Journey
    section Start
      Login: 5: User
      Upload document: 4: User
    section Ask
      Submit question: 5: User
      Review answer: 4: User
      Refine prompt: 3: User"""

    if diagram_type == "gitgraph":
        return """gitGraph
    commit id: "init"
    branch feature
    checkout feature
    commit id: "backend"
    commit id: "frontend"
    checkout main
    merge feature
    commit id: "release\""""

    if diagram_type == "quadrant":
        return """quadrantChart
    title Feature Priority
    x-axis Low Effort --> High Effort
    y-axis Low Impact --> High Impact
    quadrant-1 Strategic
    quadrant-2 Quick Wins
    quadrant-3 Avoid
    quadrant-4 Major Projects
    MCP Integration: [0.55, 0.82]
    Mermaid Rendering: [0.35, 0.75]
    File Uploads: [0.45, 0.70]"""

    if diagram_type == "requirement":
        return """requirementDiagram
    requirement diagrams {
      id: 1
      text: Render clean Mermaid diagrams
      risk: medium
      verifymethod: test
    }
    functionalRequirement renderer {
      id: 1.1
      text: Convert Mermaid blocks to SVG
      risk: low
      verifymethod: inspection
    }
    diagrams - contains -> renderer"""

    if diagram_type == "sankey":
        return """sankey-beta
    User,Frontend,10
    Frontend,Backend,10
    Backend,Retriever,6
    Backend,Model,4
    Retriever,Answer,6
    Model,Answer,4"""

    if diagram_type == "xy":
        return """xyChart-beta
    title "Requests Over Time"
    x-axis [Mon, Tue, Wed, Thu, Fri]
    y-axis "Requests" 0 --> 100
    line [20, 35, 48, 70, 85]"""

    return """flowchart TD
    A[Start] --> B[Process]
    B --> C[Finish]"""


def fallback_mermaid_diagram(user_message: str) -> str:
    requested_types = requested_mermaid_types(user_message)
    requested_count = mermaid_request_count(user_message)
    if len(requested_types) < requested_count:
        requested_types.extend(["flowchart"] * (requested_count - len(requested_types)))

    blocks = []
    for diagram_type in requested_types[:requested_count]:
        title = mermaid_title_for(user_message, diagram_type)
        heading = f"**{title} - {diagram_type.title()}**" if requested_count > 1 else f"**{title}**"
        blocks.append(f"{heading}\n```mermaid\n{mermaid_code_for(user_message, diagram_type)}\n```")

    intro = "Here are the Mermaid diagrams:" if len(blocks) > 1 else "Here is the Mermaid diagram:"
    return intro + "\n\n" + "\n\n".join(blocks)


def select_openrouter_tools(user_message: str) -> List[dict]:
    tools = []
    if wants_web_search_tool(user_message):
        tools.append(WEB_SEARCH_PLUGIN)
    return tools


def load_mcp_server_configs() -> List[dict]:
    raw_config = MCP_SERVERS_JSON.strip()
    if not raw_config and os.path.exists(MCP_SERVERS_FILE):
        try:
            with open(MCP_SERVERS_FILE, "r", encoding="utf-8") as config_file:
                raw_config = config_file.read().strip()
        except OSError as exc:
            print(f"Could not read MCP config file: {exc}")
            return []

    if not raw_config:
        return []

    try:
        parsed = json.loads(raw_config)
    except json.JSONDecodeError as exc:
        print(f"Invalid MCP server config: {exc}")
        return []

    if isinstance(parsed, dict):
        parsed = parsed.get("servers", [])

    if not isinstance(parsed, list):
        print("Invalid MCP server config: expected a list or {'servers': [...]}.")
        return []

    configs = []
    for server in parsed:
        if not isinstance(server, dict):
            continue
        name = str(server.get("name", "")).strip()
        command = str(server.get("command", "")).strip()
        args = server.get("args", [])
        env = server.get("env", {})
        if not name or not command:
            continue
        if not isinstance(args, list):
            args = []
        if not isinstance(env, dict):
            env = {}
        configs.append({
            "name": name,
            "command": command,
            "args": [str(arg) for arg in args],
            "env": {str(key): str(value) for key, value in env.items()},
        })
    return configs


def mcp_tool_name(server_name: str, tool_name: str) -> str:
    safe_server = re.sub(r"[^a-zA-Z0-9_-]+", "_", server_name).strip("_")
    safe_tool = re.sub(r"[^a-zA-Z0-9_-]+", "_", tool_name).strip("_")
    return f"mcp__{safe_server}__{safe_tool}"[:64]


def split_mcp_tool_name(name: str) -> tuple[str, str] | None:
    parts = name.split("__", 2)
    if len(parts) != 3 or parts[0] != "mcp":
        return None
    return parts[1], parts[2]


def mcp_content_to_text(content: Any) -> str:
    if isinstance(content, str):
        return content
    if isinstance(content, list):
        return "\n".join(mcp_content_to_text(item) for item in content if item is not None)
    if hasattr(content, "text"):
        return str(content.text)
    if hasattr(content, "model_dump"):
        return json.dumps(content.model_dump(), ensure_ascii=False)
    if isinstance(content, dict):
        return json.dumps(content, ensure_ascii=False)
    return str(content)


async def list_mcp_tools_async() -> List[dict]:
    configs = load_mcp_server_configs()
    if not configs:
        return []

    try:
        from mcp import ClientSession, StdioServerParameters
        from mcp.client.stdio import stdio_client
    except ImportError as exc:
        raise HTTPException(
            status_code=503,
            detail="MCP Python package is not installed. Run pip install -r backend/requirements.txt.",
        ) from exc

    openrouter_tools = []
    for config in configs:
        params = StdioServerParameters(
            command=config["command"],
            args=config["args"],
            env={**os.environ, **config["env"]},
        )
        try:
            async with stdio_client(params) as (read_stream, write_stream):
                async with ClientSession(read_stream, write_stream) as session:
                    await session.initialize()
                    result = await session.list_tools()
        except Exception as exc:
            print(f"Could not list MCP tools for {config['name']}: {exc}")
            continue

        for tool in result.tools:
            if tool.name not in MCP_ALLOWED_TOOLS:
                continue
            schema = getattr(tool, "inputSchema", None) or {"type": "object", "properties": {}}
            openrouter_tools.append({
                "type": "function",
                "function": {
                    "name": mcp_tool_name(config["name"], tool.name),
                    "description": (
                        f"MCP server '{config['name']}' tool '{tool.name}'. "
                        f"{getattr(tool, 'description', '') or ''}"
                    ).strip(),
                    "parameters": schema,
                },
            })
    return openrouter_tools


async def call_mcp_tool_async(tool_call: dict) -> str:
    function = tool_call.get("function") or {}
    mapped = split_mcp_tool_name(function.get("name", ""))
    if not mapped:
        return "Unknown MCP tool."

    safe_server_name, safe_tool_name = mapped
    configs = load_mcp_server_configs()
    config = next(
        (
            item for item in configs
            if re.sub(r"[^a-zA-Z0-9_-]+", "_", item["name"]).strip("_") == safe_server_name
        ),
        None,
    )
    if not config:
        return f"MCP server '{safe_server_name}' is not configured."

    try:
        arguments = json.loads(function.get("arguments") or "{}")
    except json.JSONDecodeError:
        arguments = {}

    try:
        from mcp import ClientSession, StdioServerParameters
        from mcp.client.stdio import stdio_client
    except ImportError as exc:
        raise HTTPException(
            status_code=503,
            detail="MCP Python package is not installed. Run pip install -r backend/requirements.txt.",
        ) from exc

    params = StdioServerParameters(
        command=config["command"],
        args=config["args"],
        env={**os.environ, **config["env"]},
    )
    async with stdio_client(params) as (read_stream, write_stream):
        async with ClientSession(read_stream, write_stream) as session:
            await session.initialize()
            tools = await session.list_tools()
            tool = next(
                (
                    item for item in tools.tools
                    if item.name in MCP_ALLOWED_TOOLS
                    and (
                        mcp_tool_name(config["name"], item.name) == function.get("name")
                        or re.sub(r"[^a-zA-Z0-9_-]+", "_", item.name).strip("_") == safe_tool_name
                    )
                ),
                None,
            )
            if not tool:
                return f"MCP tool '{safe_tool_name}' was not found on server '{config['name']}'."
            result = await session.call_tool(tool.name, arguments)

    if getattr(result, "isError", False):
        return f"MCP tool returned an error: {mcp_content_to_text(getattr(result, 'content', ''))}"
    structured = getattr(result, "structuredContent", None)
    if structured:
        return json.dumps(structured, ensure_ascii=False)
    return mcp_content_to_text(getattr(result, "content", result))


def list_mcp_tools() -> List[dict]:
    return asyncio.run(list_mcp_tools_async())


def call_mcp_tool(tool_call: dict) -> str:
    return asyncio.run(call_mcp_tool_async(tool_call))


def parse_text_tool_calls(answer: str) -> List[dict]:
    match = re.search(r"<TOOLCALL>\s*(.*?)\s*</TOOLCALL>", answer, flags=re.DOTALL)
    if not match:
        xml_match = re.search(
            r"<tool_call>\s*<function=([^>]+)>\s*(.*?)\s*</function>\s*</tool_call>",
            answer,
            flags=re.DOTALL | re.IGNORECASE,
        )
        if not xml_match:
            simple_match = re.search(
                r"<tool_call>\s*([a-zA-Z0-9_.:-]+)\s*(.*?)\s*</tool_call>",
                answer,
                flags=re.DOTALL | re.IGNORECASE,
            )
            if simple_match:
                name = simple_match.group(1).strip()
                body = simple_match.group(2)
                keys = re.findall(
                    r"<arg_key>\s*(.*?)\s*</arg_key>",
                    body,
                    flags=re.DOTALL | re.IGNORECASE,
                )
                values = re.findall(
                    r"<arg_value>\s*(.*?)\s*</arg_value>",
                    body,
                    flags=re.DOTALL | re.IGNORECASE,
                )
                arguments = {
                    key.strip(): value.strip()
                    for key, value in zip(keys, values)
                }
                return [{
                    "id": "text-tool-call-0",
                    "type": "function",
                    "function": {
                        "name": name,
                        "arguments": json.dumps(arguments),
                    },
                }]

        if not xml_match:
            minimax_match = re.search(
                r"<minimax:tool_call>\s*<invoke\s+name=[\"']([^\"']+)[\"']>\s*(.*?)\s*</invoke>\s*</minimax:tool_call>",
                answer,
                flags=re.DOTALL | re.IGNORECASE,
            )
            if minimax_match:
                name = minimax_match.group(1).strip()
                body = minimax_match.group(2)
                arguments = {}
                for param_name, param_value in re.findall(
                    r"<parameter\s+name=[\"']([^\"']+)[\"']>\s*(.*?)\s*</parameter>",
                    body,
                    flags=re.DOTALL | re.IGNORECASE,
                ):
                    arguments[param_name.strip()] = param_value.strip()

                return [{
                    "id": "text-tool-call-0",
                    "type": "function",
                    "function": {
                        "name": name,
                        "arguments": json.dumps(arguments),
                    },
                }]

        if not xml_match:
            return []

        name = xml_match.group(1).strip()
        body = xml_match.group(2)
        arguments = {}
        for param_name, param_value in re.findall(
            r"<parameter=([^>]+)>\s*(.*?)\s*</parameter>",
            body,
            flags=re.DOTALL | re.IGNORECASE,
        ):
            arguments[param_name.strip()] = param_value.strip()

        return [{
            "id": "text-tool-call-0",
            "type": "function",
            "function": {
                "name": name,
                "arguments": json.dumps(arguments),
            },
        }]

    try:
        parsed = json.loads(match.group(1))
    except json.JSONDecodeError:
        return []

    if isinstance(parsed, dict):
        parsed = [parsed]
    if not isinstance(parsed, list):
        return []

    tool_calls = []
    for index, item in enumerate(parsed[:4]):
        if not isinstance(item, dict):
            continue
        name = item.get("name")
        arguments = item.get("arguments") or {}
        if not name:
            continue
        tool_calls.append({
            "id": f"text-tool-call-{index}",
            "type": "function",
            "function": {
                "name": str(name),
                "arguments": json.dumps(arguments),
            },
        })
    return tool_calls


def resolve_mcp_tool_call(tool_call: dict, available_tools: List[dict]) -> dict:
    function = tool_call.get("function") or {}
    requested_name = function.get("name", "")
    available_names = [
        (tool.get("function") or {}).get("name", "")
        for tool in available_tools
    ]
    if requested_name in available_names:
        return tool_call

    requested_tail = requested_name.split("__")[-1]
    aliases = {
        "list_directory_contents": "list_directory",
        "list_files": "list_directory",
        "read": "read_text_file",
    }
    target_tail = aliases.get(requested_tail, requested_tail)
    replacement = next(
        (
            name for name in available_names
            if name.endswith(f"__{target_tail}")
        ),
        None,
    )
    if not replacement:
        return tool_call

    return {
        **tool_call,
        "function": {
            **function,
            "name": replacement,
        },
    }


def normalize_mcp_arguments(tool_call: dict) -> dict:
    function = tool_call.get("function") or {}
    try:
        arguments = json.loads(function.get("arguments") or "{}")
    except json.JSONDecodeError:
        arguments = {}

    def normalize_path(value: str) -> str:
        if not value or os.path.isabs(value):
            return value
        return os.path.abspath(os.path.join(WORKSPACE_DIR, value))

    if isinstance(arguments.get("path"), str):
        arguments["path"] = normalize_path(arguments["path"])
    if isinstance(arguments.get("paths"), list):
        arguments["paths"] = [
            normalize_path(path) if isinstance(path, str) else path
            for path in arguments["paths"]
        ]

    return {
        **tool_call,
        "function": {
            **function,
            "arguments": json.dumps(arguments),
        },
    }


def direct_mcp_answer(user_message: str) -> str | None:
    lowered = user_message.lower()
    if "mcp" not in lowered:
        return None

    def run_tool(name: str, arguments: dict) -> str:
        tool_call = {
            "id": "direct-mcp-call",
            "type": "function",
            "function": {
                "name": name,
                "arguments": json.dumps(arguments),
            },
        }
        return call_mcp_tool(normalize_mcp_arguments(tool_call))

    def unwrap_mcp_content(result: str) -> str:
        try:
            parsed = json.loads(result)
            if isinstance(parsed, dict) and isinstance(parsed.get("content"), str):
                return parsed["content"]
        except json.JSONDecodeError:
            pass
        return result

    def format_directory_result(result: str, label: str) -> str:
        content = unwrap_mcp_content(result)
        rows = []
        for line in content.splitlines():
            match = re.match(r"\[(FILE|DIR)\]\s+(.+)", line.strip())
            if not match:
                continue
            item_type = "File" if match.group(1) == "FILE" else "Directory"
            rows.append((item_type, match.group(2)))

        if not rows:
            return clean_llm_answer(f"MCP filesystem result for `{label}`:\n\n{content}")

        table = ["| Type | Name |", "| --- | --- |"]
        table.extend(f"| {item_type} | `{name}` |" for item_type, name in rows)
        return "\n".join([
            f"MCP filesystem result for `{label}`:",
            "",
            *table,
        ])

    def format_text_result(result: str, label: str) -> str:
        content = unwrap_mcp_content(result)
        return clean_llm_answer(f"MCP filesystem result for `{label}`:\n\n```text\n{content}\n```")

    path_match = re.search(
        r"(?:read|open|show|list|files?\s+in|directory\s+of)\s+([a-zA-Z0-9_./\\ -]+)",
        user_message,
        flags=re.IGNORECASE,
    )
    path = "backend"
    if path_match:
        path = path_match.group(1).strip(" .")
        path = path.replace(" folder", "").replace(" directory", "").strip()

    if "backend" in lowered and ("folder" in lowered or "directory" in lowered):
        path = "backend"

    if any(word in lowered for word in ("list", "files", "folder", "directory")):
        result = run_tool("mcp__workspace__list_directory", {"path": path})
        return format_directory_result(result, path)

    if any(word in lowered for word in ("read", "open", "show")):
        file_match = re.search(r"([a-zA-Z0-9_./\\ -]+\.[a-zA-Z0-9_]+)", user_message)
        file_path = file_match.group(1).strip() if file_match else path
        result = run_tool("mcp__workspace__read_text_file", {"path": file_path, "head": 120})
        return format_text_result(result, file_path)

    search_match = re.search(
        r"search(?:\s+(?:for|my workspace for))?\s+(.+?)(?:\s+in\s+([a-zA-Z0-9_./\\ -]+))?$",
        user_message,
        flags=re.IGNORECASE,
    )
    if search_match:
        query = search_match.group(1).strip(" .'\"")
        search_path = (search_match.group(2) or ".").strip(" .'\"")
        result = run_tool("mcp__workspace__search_files", {"path": search_path, "pattern": query})
        return format_text_result(result, f"search `{query}`")

    return None


def get_current_datetime_tool() -> str:
    if OPENROUTER_DATETIME_TIMEZONE.lower() in ("asia/kolkata", "asia/calcutta", "ist"):
        tz = timezone(timedelta(hours=5, minutes=30), name="IST")
    else:
        tz = timezone.utc
    now = datetime.now(tz)
    return now.strftime("%A, %B %d, %Y at %I:%M %p %Z")


def clean_llm_answer(answer: str) -> str:
    answer = answer.strip()
    answer = unicodedata.normalize("NFKC", answer)
    answer = answer.translate(str.maketrans({
        "\u2018": "'",
        "\u2019": "'",
        "\u201c": '"',
        "\u201d": '"',
        "\u2013": "-",
        "\u2014": "-",
        "\u202f": " ",
        "\u00a0": " ",
    }))
    answer = re.sub(r"【[^】]*】", "", answer)
    answer = re.sub(r"ã.*?ã", "", answer)
    answer = re.sub(r"[ \t]+\n", "\n", answer)
    answer = re.sub(r"\n{3,}", "\n\n", answer)
    return answer.strip()


def call_openrouter(
    messages: List[dict],
    plugins: List[dict] | None = None,
    tools: List[dict] | None = None,
) -> str:
    payload = {
        "model": OPENROUTER_MODEL,
        "messages": messages,
        "max_tokens": 1024,
        "temperature": 0.3,
    }
    if plugins:
        payload["plugins"] = plugins
    if tools:
        payload["tools"] = tools
        payload["tool_choice"] = "auto"

    try:
        response = requests.post(
            OPENROUTER_CHAT_URL,
            headers=openrouter_headers,
            data=json.dumps(payload),
            timeout=60,
        )
    except requests.RequestException as exc:
        raise HTTPException(status_code=502, detail=f"OpenRouter request failed: {exc}") from exc

    if response.status_code >= 400:
        error_detail = response.text[:500]
        try:
            error_payload = response.json()
            error_detail = error_payload.get("error", {}).get("message", error_detail)
        except ValueError:
            pass

        if response.status_code == 403 and "limit exceeded" in error_detail.lower():
            raise HTTPException(
                status_code=402,
                detail=(
                    "OpenRouter key limit exceeded. Increase or reset the key limit "
                    "in OpenRouter settings, then try again."
                ),
            )

        raise HTTPException(
            status_code=502,
            detail=f"OpenRouter error {response.status_code}: {error_detail}",
        )

    data = response.json()
    try:
        message = data["choices"][0]["message"]
        content = message.get("content") or ""
        tool_calls = message.get("tool_calls") or parse_text_tool_calls(content)
        if tool_calls:
            follow_up_messages = [
                *messages,
                {
                    "role": "assistant",
                    "content": message.get("content") or "",
                    "tool_calls": tool_calls,
                },
            ]
            tool_outputs = []
            for tool_call in tool_calls[:4]:
                resolved_tool_call = normalize_mcp_arguments(
                    resolve_mcp_tool_call(tool_call, tools or [])
                )
                tool_output = call_mcp_tool(resolved_tool_call)[:8000]
                tool_outputs.append(tool_output)
                follow_up_messages.append({
                    "role": "tool",
                    "tool_call_id": resolved_tool_call.get("id"),
                    "name": (resolved_tool_call.get("function") or {}).get("name"),
                    "content": tool_output,
                })
            follow_answer = call_openrouter(follow_up_messages, plugins=plugins)
            if follow_answer:
                return follow_answer
            return clean_llm_answer("MCP tool result:\n\n" + "\n\n".join(tool_outputs))
        parsed_late_tool_calls = parse_text_tool_calls(content)
        if parsed_late_tool_calls and tools is None:
            follow_up_messages = [
                *messages,
                {
                    "role": "assistant",
                    "content": content,
                    "tool_calls": parsed_late_tool_calls,
                },
            ]
            tool_outputs = []
            for tool_call in parsed_late_tool_calls[:4]:
                resolved_tool_call = normalize_mcp_arguments(
                    resolve_mcp_tool_call(tool_call, list_mcp_tools())
                )
                tool_output = call_mcp_tool(resolved_tool_call)[:8000]
                tool_outputs.append(tool_output)
                follow_up_messages.append({
                    "role": "tool",
                    "tool_call_id": resolved_tool_call.get("id"),
                    "name": (resolved_tool_call.get("function") or {}).get("name"),
                    "content": tool_output,
                })
            follow_answer = call_openrouter(follow_up_messages, plugins=plugins)
            if follow_answer:
                return follow_answer
            return clean_llm_answer("MCP tool result:\n\n" + "\n\n".join(tool_outputs))
        return clean_llm_answer(content)
    except (KeyError, IndexError, TypeError) as exc:
        raise HTTPException(status_code=502, detail="Unexpected OpenRouter response.") from exc

# ── Supabase client ───────────────────────────────────────────────
supabase: Client = create_client(SUPABASE_URL, SUPABASE_KEY)

# ── Embedding model ───────────────────────────────────────────────
print("Loading embedding model...")
embedder = SentenceTransformer("all-MiniLM-L6-v2")
EMBED_DIM = 384

# ── FAISS index ───────────────────────────────────────────────────
index = faiss.IndexFlatIP(EMBED_DIM)
metadata_store: List[dict] = []

# ── Active sessions ───────────────────────────────────────────────
active_sessions = {}

DEV_USERS = {
    "EMB001": "123456",
    "EMB002": "654321",
    "EMB003": "111222",
}

# ── FastAPI app ───────────────────────────────────────────────────
app = FastAPI(title="EMB RAG Chatbot")

app.add_middleware(
    CORSMiddleware,
    allow_origins=FRONTEND_ORIGINS,
    allow_origin_regex=r"https://.*\.pages\.dev|http://localhost:\d+|http://127\.0\.0\.1:\d+",
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

security = HTTPBearer()

# ── Request/Response models ───────────────────────────────────────
class LoginRequest(BaseModel):
    user_id: str
    passcode: str

class SignupRequest(BaseModel):
    user_id: str
    passcode: str

class UserLookupRequest(BaseModel):
    user_id: str

class ResetPasswordRequest(BaseModel):
    user_id: str
    new_passcode: str

class LoginResponse(BaseModel):
    success: bool
    token: str
    message: str

class Message(BaseModel):
    role: str
    content: str

class ChatRequest(BaseModel):
    message: str
    history: List[Message] = Field(default_factory=list)
    top_k: int = 4

class ChatResponse(BaseModel):
    answer: str
    sources: List[str]
    chunks_used: int

class MemoryRequest(BaseModel):
    sessions: List[dict] = Field(default_factory=list)
    active_session_id: str = ""
    history: List[Message] = Field(default_factory=list)

class MemoryResponse(BaseModel):
    sessions: List[dict]
    active_session_id: str
    history: List[Message]


# ================================================================
#  AUTH FUNCTIONS
# ================================================================

def verify_token(credentials: HTTPAuthorizationCredentials = Depends(security)):
    token = credentials.credentials
    if token not in active_sessions:
        raise HTTPException(
            status_code=401,
            detail="Invalid or expired session. Please login again."
        )
    return active_sessions[token]


def normalize_auth_identifier(value: str) -> str:
    identifier = value.strip()
    if "@" in identifier:
        return identifier.lower()
    return identifier.upper()


def get_supabase_users(user_id: str, columns: str = "*") -> List[dict]:
    try:
        result = supabase.table("users").select(columns).eq(
            "user_id", user_id
        ).execute()

        if not result.data:
            result = supabase.table("users").select(columns).ilike(
                "user_id", user_id
            ).execute()

        return result.data or []
    except Exception as e:
        print(f"Could not check Supabase users: {e}")
        return []


def user_exists(user_id: str) -> bool:
    return bool(get_supabase_users(user_id, "user_id")) or user_id in DEV_USERS


# ================================================================
#  MEMORY FUNCTIONS
# ================================================================

MEMORY_TABLE = "chat_memories"


def memory_table_error(exc: Exception) -> HTTPException:
    return HTTPException(
        status_code=503,
        detail=(
            "Chat memory table is not ready. Create a Supabase table named "
            "chat_memories with user_id, sessions, active_session_id, history, "
            "and updated_at columns."
        ),
    )


def load_user_memory(user_id: str) -> dict:
    try:
        result = supabase.table(MEMORY_TABLE).select(
            "sessions, active_session_id, history"
        ).eq("user_id", user_id).execute()
    except Exception as exc:
        print(f"Could not load chat memory: {exc}")
        raise memory_table_error(exc) from exc

    if not result.data:
        return {"sessions": [], "active_session_id": "", "history": []}

    row = result.data[0]
    return {
        "sessions": row.get("sessions") or [],
        "active_session_id": row.get("active_session_id") or "",
        "history": row.get("history") or [],
    }


def save_user_memory(user_id: str, memory: MemoryRequest) -> None:
    try:
        supabase.table(MEMORY_TABLE).upsert({
            "user_id": user_id,
            "sessions": memory.sessions,
            "active_session_id": memory.active_session_id,
            "history": [msg.model_dump() for msg in memory.history],
            "updated_at": datetime.now(timezone.utc).isoformat(),
        }, on_conflict="user_id").execute()
    except Exception as exc:
        print(f"Could not save chat memory: {exc}")
        raise memory_table_error(exc) from exc


# ================================================================
#  RAG HELPER FUNCTIONS
# ================================================================

def extract_text_from_pdf(file_bytes: bytes) -> str:
    try:
        from pypdf import PdfReader
        pdf_file = io.BytesIO(file_bytes)
        reader = PdfReader(pdf_file)
        all_text = []
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text and page_text.strip():
                all_text.append(page_text)
        full_text = "\n\n".join(all_text)
        if not full_text.strip():
            raise HTTPException(
                status_code=400,
                detail="Could not extract text. PDF may be image-based."
            )
        return full_text
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(
            status_code=400,
            detail=f"Failed to read PDF: {str(e)}"
        )


def extract_text_from_docx(file_bytes: bytes) -> str:
    try:
        from docx import Document
        document = Document(io.BytesIO(file_bytes))
        parts = []
        for paragraph in document.paragraphs:
            text = paragraph.text.strip()
            if text:
                parts.append(text)
        for table in document.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    parts.append(" | ".join(cells))
        content = "\n".join(parts)
        if not content.strip():
            raise HTTPException(status_code=400, detail="No readable text found in Word document.")
        return content
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Failed to read Word document: {str(e)}")


def extract_text_from_csv(file_bytes: bytes) -> str:
    try:
        import csv
        decoded = file_bytes.decode("utf-8-sig", errors="ignore")
        rows = []
        reader = csv.reader(io.StringIO(decoded))
        for row in reader:
            cleaned = [cell.strip() for cell in row if cell.strip()]
            if cleaned:
                rows.append(" | ".join(cleaned))
        content = "\n".join(rows)
        if not content.strip():
            raise HTTPException(status_code=400, detail="No readable text found in CSV file.")
        return content
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Failed to read CSV file: {str(e)}")


def extract_text_from_xlsx(file_bytes: bytes) -> str:
    try:
        from openpyxl import load_workbook
        workbook = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
        parts = []
        for sheet in workbook.worksheets:
            parts.append(f"Sheet: {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                values = [str(value).strip() for value in row if value is not None and str(value).strip()]
                if values:
                    parts.append(" | ".join(values))
        content = "\n".join(parts)
        if not content.strip():
            raise HTTPException(status_code=400, detail="No readable text found in Excel file.")
        return content
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Failed to read Excel file: {str(e)}")


def extract_text_from_pptx(file_bytes: bytes) -> str:
    try:
        from pptx import Presentation
        presentation = Presentation(io.BytesIO(file_bytes))
        parts = []
        for slide_index, slide in enumerate(presentation.slides, 1):
            slide_parts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_parts.append(shape.text.strip())
            if slide_parts:
                parts.append(f"Slide {slide_index}:\n" + "\n".join(slide_parts))
        content = "\n\n".join(parts)
        if not content.strip():
            raise HTTPException(status_code=400, detail="No readable text found in PowerPoint file.")
        return content
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Failed to read PowerPoint file: {str(e)}")


def extract_text_from_image(file_bytes: bytes) -> str:
    TesseractNotFoundError = RuntimeError
    try:
        from PIL import Image
        import pytesseract
        from pytesseract import TesseractNotFoundError
        if TESSERACT_CMD and os.path.exists(TESSERACT_CMD):
            pytesseract.pytesseract.tesseract_cmd = TESSERACT_CMD
        image = Image.open(io.BytesIO(file_bytes))
        content = pytesseract.image_to_string(image)
        if not content.strip():
            raise HTTPException(status_code=400, detail="No readable text found in image.")
        return content
    except HTTPException:
        raise
    except TesseractNotFoundError as exc:
        raise HTTPException(
            status_code=400,
            detail=(
                "Image OCR needs Tesseract installed on Windows. Install Tesseract OCR "
                "or upload text-based documents instead."
            ),
        ) from exc
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Failed to read image text: {str(e)}")


def chunk_text(text: str, chunk_size: int = 300, overlap: int = 50) -> List[str]:
    words = text.split()
    chunks = []
    i = 0
    while i < len(words):
        chunk = " ".join(words[i: i + chunk_size])
        chunks.append(chunk)
        i += chunk_size - overlap
    return [c for c in chunks if len(c.strip()) > 30]


def embed(texts: List[str]) -> np.ndarray:
    vectors = embedder.encode(texts, convert_to_numpy=True)
    norms = np.linalg.norm(vectors, axis=1, keepdims=True)
    normalized = vectors / (norms + 1e-9)
    return normalized.astype("float32")


def retrieve(query: str, top_k: int = 4) -> List[dict]:
    if index.ntotal == 0:
        return []
    q_vec = embed([query])
    scores, positions = index.search(q_vec, min(top_k, index.ntotal))
    results = []
    for score, pos in zip(scores[0], positions[0]):
        if pos == -1:
            continue
        results.append({
            "text": metadata_store[pos]["text"],
            "source": metadata_store[pos]["source"],
            "score": float(score),
        })
    return results


def build_rag_prompt(query: str, chunks: List[dict]) -> str:
    context_parts = []
    for i, chunk in enumerate(chunks, 1):
        context_parts.append(
            f"[Source {i} — {chunk['source']}]\n{chunk['text']}"
        )
    context = "\n\n---\n\n".join(context_parts)
    return f"""You are a helpful assistant. Answer ONLY using the context below.
If the answer is not in the context, say: "I couldn't find that in the uploaded documents."
Always cite sources as [Source 1], [Source 2], etc.

CONTEXT:
{context}

USER QUESTION:
{query}"""


# ================================================================
#  API ROUTES
# ================================================================

@app.get("/health")
def health():
    return {
        "status": "ok",
        "chunks_indexed": index.ntotal,
        "model": f"{OPENROUTER_MODEL} (OpenRouter)",
    }


@app.post("/signup")
def signup(req: SignupRequest):
    try:
        user_id = normalize_auth_identifier(req.user_id)
        passcode = req.passcode.strip()

        if not user_id:
            raise HTTPException(status_code=400, detail="Email address is required.")
        if len(passcode) != 6 or not passcode.isdigit():
            raise HTTPException(status_code=400, detail="Passcode must be 6 digits.")

        if user_exists(user_id):
            raise HTTPException(status_code=409, detail="Account already exists.")

        try:
            supabase.table("users").insert({
                "user_id": user_id,
                "passcode": passcode,
                "created_at": datetime.now(timezone.utc).isoformat(),
            }).execute()
        except Exception as e:
            print(f"Could not save signup to Supabase, using local dev user: {e}")

        DEV_USERS[user_id] = passcode

        return {
            "success": True,
            "message": "Signup successful. Please login.",
            "user_id": user_id,
        }

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/forgot-password")
def forgot_password(req: UserLookupRequest):
    try:
        user_id = normalize_auth_identifier(req.user_id)

        if not user_id:
            raise HTTPException(status_code=400, detail="Email address is required.")

        if not user_exists(user_id):
            raise HTTPException(status_code=404, detail="Account not found.")

        return {
            "success": True,
            "message": "User found. You can reset your passcode now.",
            "user_id": user_id,
        }

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/reset-password")
def reset_password(req: ResetPasswordRequest):
    try:
        user_id = normalize_auth_identifier(req.user_id)
        new_passcode = req.new_passcode.strip()

        if not user_id:
            raise HTTPException(status_code=400, detail="Email address is required.")
        if len(new_passcode) != 6 or not new_passcode.isdigit():
            raise HTTPException(status_code=400, detail="New passcode must be 6 digits.")
        if not user_exists(user_id):
            raise HTTPException(status_code=404, detail="Account not found.")

        supabase_users = get_supabase_users(user_id, "user_id")
        if supabase_users:
            try:
                supabase.table("users").update({
                    "passcode": new_passcode,
                    "updated_at": datetime.now(timezone.utc).isoformat(),
                }).eq("user_id", supabase_users[0]["user_id"]).execute()
            except Exception as e:
                print(f"Could not reset passcode in Supabase, using local dev user: {e}")

        DEV_USERS[user_id] = new_passcode

        return {
            "success": True,
            "message": "Passcode reset successful. Please login.",
            "user_id": user_id,
        }

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/login", response_model=LoginResponse)
def login(req: LoginRequest):
    try:
        user_id = normalize_auth_identifier(req.user_id)
        passcode = req.passcode.strip()

        users = get_supabase_users(user_id)

        if users:
            user = users[0]

            if user["passcode"] != passcode:
                raise HTTPException(
                    status_code=401,
                    detail="Incorrect passcode."
                )

            try:
                supabase.table("users").update(
                    {"last_login": datetime.now(timezone.utc).isoformat()}
                ).eq("user_id", user["user_id"]).execute()
            except Exception as e:
                print(f"Could not update last_login: {e}")
        else:
            if DEV_USERS.get(user_id) != passcode:
                raise HTTPException(
                    status_code=401,
                    detail="Account not found or incorrect passcode."
                )

        token = secrets.token_hex(32)
        active_sessions[token] = user_id

        return LoginResponse(
            success=True,
            token=token,
            message=f"Welcome {user_id}!"
        )

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/logout")
def logout(
    user_id: str = Depends(verify_token),
    credentials: HTTPAuthorizationCredentials = Depends(security)
):
    token = credentials.credentials
    if token in active_sessions:
        del active_sessions[token]
    return {"message": "Logged out successfully"}


@app.get("/memory", response_model=MemoryResponse)
def get_memory(user_id: str = Depends(verify_token)):
    return MemoryResponse(**load_user_memory(user_id))


@app.put("/memory")
def save_memory(req: MemoryRequest, user_id: str = Depends(verify_token)):
    save_user_memory(user_id, req)
    return {"message": "Chat memory saved."}


@app.delete("/memory")
def clear_memory(user_id: str = Depends(verify_token)):
    try:
        supabase.table(MEMORY_TABLE).delete().eq("user_id", user_id).execute()
    except Exception as exc:
        print(f"Could not clear chat memory: {exc}")
        raise memory_table_error(exc) from exc
    return {"message": "Chat memory cleared."}


@app.get("/mcp/tools")
def get_mcp_tools(user_id: str = Depends(verify_token)):
    configured_servers = [server["name"] for server in load_mcp_server_configs()]
    return {
        "servers": configured_servers,
        "tools": list_mcp_tools(),
    }


@app.post("/ingest/file")
async def ingest_file(
    file: UploadFile = File(...),
    user_id: str = Depends(verify_token)
):
    filename = file.filename
    file_bytes = await file.read()

    if filename.lower().endswith(".pdf"):
        content = extract_text_from_pdf(file_bytes)
    elif filename.lower().endswith((".txt", ".md")):
        content = file_bytes.decode("utf-8", errors="ignore")
    elif filename.lower().endswith(".docx"):
        content = extract_text_from_docx(file_bytes)
    elif filename.lower().endswith(".csv"):
        content = extract_text_from_csv(file_bytes)
    elif filename.lower().endswith(".xlsx"):
        content = extract_text_from_xlsx(file_bytes)
    elif filename.lower().endswith(".pptx"):
        content = extract_text_from_pptx(file_bytes)
    elif filename.lower().endswith((".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tiff", ".tif")):
        content = extract_text_from_image(file_bytes)
    else:
        raise HTTPException(
            status_code=400,
            detail=(
                "Supported files: .pdf, .txt, .md, .docx, .csv, .xlsx, "
                ".pptx, .png, .jpg, .jpeg, .webp, .bmp, .tiff, .tif."
            )
        )

    chunks = chunk_text(content)
    if not chunks:
        raise HTTPException(
            status_code=400,
            detail="No text found in file."
        )

    vectors = embed(chunks)
    index.add(vectors)

    for chunk in chunks:
        metadata_store.append({
            "text": chunk,
            "source": filename,
        })

    return {
        "message": f"Ingested {len(chunks)} chunks from '{filename}'",
        "total_indexed": index.ntotal,
    }


@app.post("/chat", response_model=ChatResponse)
def chat(
    req: ChatRequest,
    user_id: str = Depends(verify_token)
):
    direct_answer = direct_mcp_answer(req.message)
    if direct_answer:
        return ChatResponse(
            answer=direct_answer,
            sources=[],
            chunks_used=0,
        )

    if wants_mermaid_diagram(req.message):
        return ChatResponse(
            answer=fallback_mermaid_diagram(req.message),
            sources=[],
            chunks_used=0,
        )

    tools = select_openrouter_tools(req.message)
    use_web_search = bool(tools)
    diagram_requested = wants_mermaid_diagram(req.message)
    mcp_tools = [] if use_web_search or diagram_requested else list_mcp_tools()
    chunks = [] if use_web_search else retrieve(req.message, top_k=req.top_k)
    sources = [] if use_web_search else list({c["source"] for c in chunks})

    if wants_datetime_tool(req.message) and not tools and not chunks:
        return ChatResponse(
            answer=f"The current date and time is **{get_current_datetime_tool()}**.",
            sources=[],
            chunks_used=0,
        )

    messages = []
    messages.append({
        "role": "system",
        "content": (
            "You are a helpful AI assistant. "
            "Use the conversation history to understand follow-up questions, references, and user preferences. "
            "Answer questions based on provided document context. "
            "If no document context is available, answer normally using conversation history and general knowledge. "
            "If the user asks for web, live, latest, current, recent, news, online, or search results, "
            "ignore uploaded document context and use web search tool context directly. "
            "If MCP tools are available and useful, call them before answering. "
            "Be concise, accurate, and use clean markdown formatting. "
            "Use plain ASCII punctuation. "
            "For web search answers, start with the direct answer, then add 2-4 short bullets. "
            "For Mermaid diagram requests, return a short title or note and fenced "
            "```mermaid code block(s). If the user asks for multiple diagrams, return exactly "
            "that many Mermaid code blocks. Generate valid Mermaid syntax only inside each block. "
            "Support Mermaid diagram types including flowchart, sequenceDiagram, classDiagram, "
            "stateDiagram-v2, erDiagram, journey, gantt, pie, quadrantChart, requirementDiagram, "
            "gitGraph, mindmap, timeline, C4Context, block-beta, packet-beta, architecture-beta, "
            "sankey-beta, and xyChart-beta. Choose the best type if the user does not specify one. "
            "Do not wrap Mermaid diagrams in JSON and do not include extra alternative diagrams. "
            "Do not include bracketed citation markers such as [1], line numbers, or source snippets. "
            "Do not dump raw tool output."
        ),
    })

    if wants_datetime_tool(req.message):
        messages.append({
            "role": "system",
            "content": f"Current date/time tool result: {get_current_datetime_tool()}.",
        })

    if wants_mermaid_diagram(req.message):
        messages.append({
            "role": "system",
            "content": (
                "The user wants a renderable Mermaid diagram. Keep labels short, quote labels "
                "that contain punctuation, avoid markdown inside node labels, and ensure the "
                "diagram code is syntactically valid Mermaid."
            ),
        })

    if mcp_tools:
        tool_names = ", ".join(
            (tool.get("function") or {}).get("name", "")
            for tool in mcp_tools
        )
        messages.append({
            "role": "system",
            "content": (
                "Available MCP tool names: "
                f"{tool_names}. Use only these exact names when calling MCP tools."
            ),
        })

    for msg in req.history[-8:]:
        messages.append({"role": msg.role, "content": msg.content})

    if use_web_search:
        user_content = (
            "Use web search for this request. Do not answer from uploaded documents unless the user "
            f"explicitly asks to compare with them.\n\nUSER REQUEST:\n{req.message}"
        )
    elif chunks:
        user_content = build_rag_prompt(req.message, chunks)
    elif wants_datetime_tool(req.message) or req.history:
        user_content = req.message
    else:
        user_content = req.message

    messages.append({"role": "user", "content": user_content})

    answer = call_openrouter(messages, plugins=tools, tools=mcp_tools)
    if not answer:
        answer = fallback_mermaid_diagram(req.message) if diagram_requested else (
            "I could not generate a response from the current OpenRouter model. Please try again."
        )

    return ChatResponse(
        answer=answer,
        sources=sources,
        chunks_used=len(chunks),
    )


@app.delete("/vectorstore")
def clear_vectorstore(user_id: str = Depends(verify_token)):
    global index, metadata_store
    index = faiss.IndexFlatIP(EMBED_DIM)
    metadata_store.clear()
    return {"message": "Vector store cleared."}
